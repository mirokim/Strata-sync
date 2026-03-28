#!/usr/bin/env python3
"""
pptx_to_md.py — Graph RAG 데이터 정제 매뉴얼 v3.8 §4.3 구현
PPTX → Obsidian 호환 Markdown 변환

사용법:
  python pptx_to_md.py <pptx_path> --active <active_dir> --attachments <att_dir>
  python pptx_to_md.py <pptx_dir>  --active <active_dir> --attachments <att_dir>
"""

import os
import re
import sys
import shutil
import argparse
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import pptx.oxml.ns as ns
except ImportError:
    print("ERROR: python-pptx 미설치. pip install python-pptx --break-system-packages")
    sys.exit(1)


def sanitize_filename(name: str) -> str:
    return re.sub(r'[\\/*?:"<>|]', '_', name).strip()


def detect_type(filename: str) -> str:
    fname = filename.lower()
    if any(kw in fname for kw in ['회의', '피드백', '보고', 'meeting']):
        return 'meeting'
    if any(kw in fname for kw in ['가이드', 'guide', '매뉴얼', '튜토리얼']):
        return 'guide'
    if any(kw in fname for kw in ['결정', 'decision', 'adr']):
        return 'decision'
    return 'spec'


def detect_tags(filename: str) -> list:
    tags = []
    for kw in ['이사장', '피드백', '정례보고', '정례 보고', '회장님']:
        if kw in filename:
            tags.append('chief')
            break
    return tags


def shape_to_text(shape) -> str:
    """Shape에서 텍스트 추출."""
    parts = []

    # 텍스트프레임
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = ''.join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)

    # 테이블
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace('|', '\\|').replace('\n', ' ')
                     for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
        if rows:
            col_count = len(table.columns)
            rows.insert(1, '|' + '---|' * col_count)
            parts.append('\n'.join(rows))

    return '\n'.join(parts)


def extract_slide_images(slide, stem: str, slide_num: int, attachments_dir: Path) -> list:
    """슬라이드에서 삽입 이미지 추출 → attachments/ 저장."""
    image_links = []
    img_counter = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_counter += 1
            try:
                image = shape.image
                ext = image.ext  # 'png', 'jpg' 등
                filename = f"{stem}_p{slide_num}_{img_counter}.{ext}"
                dst = attachments_dir / filename
                dst.write_bytes(image.blob)
                image_links.append(filename)
            except Exception:
                pass

    return image_links


def pptx_to_md(pptx_path: Path, active_dir: Path, attachments_dir: Path) -> dict:
    """단일 PPTX 파일 변환."""
    stem = sanitize_filename(pptx_path.stem)
    title = pptx_path.stem
    doc_type = detect_type(title)
    tags = detect_tags(title)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {'file': str(pptx_path), 'status': 'error', 'msg': str(e)}

    # 수정일 추출
    try:
        mtime = datetime.fromtimestamp(pptx_path.stat().st_mtime).strftime('%Y-%m-%d')
    except Exception:
        mtime = datetime.now().strftime('%Y-%m-%d')

    slides_md = []
    total_images = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        # 슬라이드 제목 추출
        slide_title = ''
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text_frame.text.strip()

        heading = f"## 슬라이드 {slide_num}" + (f" — {slide_title}" if slide_title else '')

        # 본문 텍스트 수집 (제목 shape 제외)
        body_parts = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            text = shape_to_text(shape)
            if text.strip():
                body_parts.append(text)

        # 이미지 추출
        image_links = extract_slide_images(slide, stem, slide_num, attachments_dir)
        total_images += len(image_links)

        # 발표자 노트
        notes_text = ''
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # 슬라이드 MD 조합
        slide_md_parts = [heading]
        if body_parts:
            slide_md_parts.append('\n'.join(body_parts))
        for img in image_links:
            slide_md_parts.append(f'![[{img}]]')
        if notes_text:
            slide_md_parts.append(f'> 📌 노트: {notes_text}')

        slides_md.append('\n\n'.join(slide_md_parts))

    # Frontmatter
    tags_yaml = '[' + ', '.join(tags) + ']' if tags else '[]'
    frontmatter = f"""---
title: "{title.replace('"', "'")}"
date: {mtime}
type: {doc_type}
status: active
tags: {tags_yaml}
source: "{pptx_path.name}"
origin: pptx
---"""

    body = '\n\n'.join(slides_md)
    md_content = frontmatter + '\n\n' + f"# {title}\n\n" + body

    out_path = active_dir / f"{stem}.md"
    try:
        out_path.write_text(md_content, encoding='utf-8')
    except Exception as e:
        return {'file': str(pptx_path), 'status': 'error', 'msg': str(e)}

    return {
        'file': str(pptx_path),
        'status': 'ok',
        'out': str(out_path),
        'slides': len(prs.slides),
        'images': total_images,
    }


def main():
    parser = argparse.ArgumentParser(description='PPTX → Obsidian MD 변환 (§4.3)')
    parser.add_argument('input', nargs='+', help='PPTX 파일 또는 폴더')
    parser.add_argument('--active', default='refined_vault/active', help='MD 출력 폴더')
    parser.add_argument('--attachments', default='refined_vault/attachments', help='이미지 출력 폴더')
    args = parser.parse_args()

    active_dir = Path(args.active)
    attachments_dir = Path(args.attachments)
    active_dir.mkdir(parents=True, exist_ok=True)
    attachments_dir.mkdir(parents=True, exist_ok=True)

    pptx_files = []
    for inp in args.input:
        p = Path(inp)
        if p.is_dir():
            pptx_files.extend(sorted(p.rglob('*.pptx')))
        elif p.suffix.lower() == '.pptx':
            pptx_files.append(p)

    print(f"PPTX {len(pptx_files)}개 변환 시작...")

    ok, errors = 0, 0
    for pptx_path in pptx_files:
        r = pptx_to_md(pptx_path, active_dir, attachments_dir)
        if r['status'] == 'ok':
            ok += 1
            print(f"  ✓ {pptx_path.name} ({r.get('slides',0)} 슬라이드, {r.get('images',0)} 이미지)")
        else:
            errors += 1
            print(f"  ✗ {pptx_path.name} — {r.get('msg','')}")

    print(f"\n완료: 성공 {ok}개, 오류 {errors}개")


if __name__ == '__main__':
    main()
